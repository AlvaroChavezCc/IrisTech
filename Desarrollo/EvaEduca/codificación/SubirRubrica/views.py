from django.shortcuts import render, redirect
from django.http import HttpResponse
from .forms import UploadRubricaForm
from .models import Rubrica
import os
from docx import Document
from pptx import Presentation
from PyPDF2 import PdfReader

# Create your views here.

def procesar_archivo(file):
    file_extension = os.path.splitext(file.name)[1].lower()
    if file_extension == '.pdf':
        reader = PdfReader(file)
        text = "".join([page.extract_text().replace("\n", " ") for page in reader.pages])
    elif file_extension == '.docx':
        doc = Document(file)
        text = "\n".join([para.text for para in doc.paragraphs])
    elif file_extension == '.pptx':
        presentation = Presentation(file)
        text = "\n".join([shape.text for slide in presentation.slides for shape in slide.shapes if hasattr(shape, "text")])
    elif file_extension == '.txt':
        text = file.read().decode('utf-8')
    else:
        raise ValueError(f"Formato de archivo {file_extension} no soportado.")
    return text

def upload_rubrica(request):
    if request.method == 'POST':
        form = UploadRubricaForm(request.POST, request.FILES)
        if form.is_valid():
            rubrica_file = request.FILES['file']
            rubrica_text = procesar_archivo(rubrica_file)
            Rubrica.objects.create(nombre=rubrica_file.name, contenido=rubrica_text)
            return HttpResponse('Rúbrica subida y procesada exitosamente')
    else:
        form = UploadRubricaForm()
    return render(request, 'SubirRubrica/upload_rubrica.html', {'form': form})
